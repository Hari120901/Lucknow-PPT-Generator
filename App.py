import streamlit as st
import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from google.oauth2 import service_account
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload
from PIL import Image

st.set_page_config(page_title="Poster Frames PPT", layout="wide")
st.title("📊 Poster Frames POP PPT Generator")

# -------------------------
# Google Drive Authentication
# -------------------------
def authenticate_drive():
    credentials = service_account.Credentials.from_service_account_info(
        st.secrets["gdrive"],
        scopes=["https://www.googleapis.com/auth/drive.readonly"],
    )
    service = build("drive", "v3", credentials=credentials)
    return service


# -------------------------
# Extract Folder ID
# -------------------------
def extract_folder_id(link):
    if "folders/" not in link:
        st.error("Invalid Google Drive Folder Link")
        st.stop()

    return link.split("folders/")[1].split("?")[0]


# -------------------------
# Get Subfolders
# -------------------------
def get_subfolders(service, parent_id):

    folders = []
    page_token = None

    while True:

        response = service.files().list(
            q=f"'{parent_id}' in parents and mimeType='application/vnd.google-apps.folder' and trashed=false",
            fields="nextPageToken, files(id, name)",
            pageToken=page_token
        ).execute()

        folders.extend(response.get("files", []))
        page_token = response.get("nextPageToken")

        if not page_token:
            break

    return folders


# -------------------------
# Get Images from Folder
# -------------------------
def get_images_in_folder(service, folder_id):

    images = []
    page_token = None

    while True:

        response = service.files().list(
            q=f"'{folder_id}' in parents and mimeType contains 'image/' and trashed=false",
            fields="nextPageToken, files(id, name)",
            pageToken=page_token
        ).execute()

        images.extend(response.get("files", []))
        page_token = response.get("nextPageToken")

        if not page_token:
            break

    return images


# -------------------------
# Download Image
# -------------------------
def download_image(service, file_id):

    request = service.files().get_media(fileId=file_id)
    fh = io.BytesIO()

    downloader = MediaIoBaseDownload(fh, request)

    done = False

    while not done:
        _, done = downloader.next_chunk()

    fh.seek(0)

    return fh


# -------------------------
# Optimize + Rotate Image
# -------------------------
def optimize_image(image_stream):

    img = Image.open(image_stream)

    # ✅ Rotate 90° clockwise
    img = img.rotate(-90, expand=True)

    # Resize (prevents memory issues)
    img.thumbnail((1500, 1500))

    output = io.BytesIO()
    img.convert("RGB").save(output, format="JPEG", quality=85)

    output.seek(0)

    return output


# -------------------------
# USER INPUT
# -------------------------
campaign_input = st.text_input("📌 Campaign Name")
drive_link = st.text_input("🔗 Google Drive Folder Link")

generate_btn = st.button("🚀 Generate Presentation")


# -------------------------
# GENERATE PPT
# -------------------------
if generate_btn:

    if not campaign_input or not drive_link:
        st.warning("Please fill all fields")
        st.stop()

    try:

        st.info("Connecting to Google Drive...")

        service = authenticate_drive()
        main_folder_id = extract_folder_id(drive_link)

        subfolders = get_subfolders(service, main_folder_id)

        if not subfolders:
            st.error("No subfolders found")
            st.stop()

        st.success(f"Found {len(subfolders)} folders")

        prs = Presentation()

        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        # Colors
        TEAL = RGBColor(0, 150, 160)
        GREY = RGBColor(242, 242, 242)
        FRENCH_NAVY = RGBColor(11, 35, 65)

        image_width = Inches(3)
        gap = Inches(3)
        top_position = Inches(1.8)

        left_positions = [
            Inches(1.3),
            Inches(1.3) + image_width + gap
        ]

        progress = st.progress(0)
        slide_count = 0
        total_images = 0

        for folder in subfolders:
            total_images += len(get_images_in_folder(service, folder["id"]))

        processed = 0

        for folder in subfolders:

            images = get_images_in_folder(service, folder["id"])

            if not images:
                continue

            for i in range(0, len(images), 2):

                slide = prs.slides.add_slide(prs.slide_layouts[6])
                slide_count += 1

                # Background
                bg = slide.background
                fill = bg.fill
                fill.solid()
                fill.fore_color.rgb = GREY

                # Campaign Label
                label_box = slide.shapes.add_textbox(
                    Inches(0.9),
                    Inches(0.6),
                    Inches(5),
                    Inches(0.6)
                )

                p = label_box.text_frame.paragraphs[0]
                p.text = "C A M P A I G N  N A M E:"
                p.font.size = Pt(24)
                p.font.name = "Montserrat"
                p.font.color.rgb = FRENCH_NAVY

                # Campaign Name
                name_box = slide.shapes.add_textbox(
                    Inches(0.9),
                    Inches(1.0),
                    Inches(7),
                    Inches(1)
                )

                p = name_box.text_frame.paragraphs[0]
                p.text = campaign_input
                p.font.size = Pt(32)
                p.font.name = "Montserrat"
                p.font.bold = True
                p.font.color.rgb = FRENCH_NAVY

                # Store Name
                store_box = slide.shapes.add_textbox(
                    Inches(7),
                    Inches(0.8),
                    Inches(4),
                    Inches(1)
                )

                p = store_box.text_frame.paragraphs[0]
                p.text = folder["name"].upper()
                p.font.size = Pt(28)
                p.font.name = "Montserrat"
                p.font.bold = True
                p.font.color.rgb = TEAL
                p.alignment = PP_ALIGN.CENTER

                slide_images = images[i:i+2]

                for idx, img in enumerate(slide_images):

                    try:
                        img_stream = download_image(service, img["id"])
                        optimized = optimize_image(img_stream)

                        picture = slide.shapes.add_picture(
                            optimized,
                            left_positions[idx],
                            top_position,
                            width=image_width
                        )

                        border = slide.shapes.add_shape(
                            1,
                            picture.left,
                            picture.top,
                            picture.width,
                            picture.height
                        )

                        border.fill.background()
                        border.line.color.rgb = TEAL
                        border.line.width = Pt(3)

                        processed += 1
                        progress.progress(processed / total_images)

                    except Exception:
                        st.warning(f"Skipping image: {img['name']}")

        # Save PPT
        ppt_io = io.BytesIO()
        prs.save(ppt_io)
        ppt_io.seek(0)

        st.success(f"✅ Presentation Generated ({slide_count} slides)")

        st.download_button(
            label="📥 Download PPT",
            data=ppt_io,
            file_name=f"{campaign_input}_Report.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )

    except Exception as e:
        st.error(f"Error: {str(e)}")
